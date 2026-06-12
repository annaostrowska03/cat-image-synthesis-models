"""Convert a PDF presentation to PPTX (one image per slide).

Each slide is rendered at high DPI and embedded as a full-bleed image,
so the PPTX looks identical to the PDF and can be edited in
PowerPoint / Canva by adding shapes/text on top.

Usage:
    python scripts/pdf_to_pptx.py presentation/main.pdf presentation/main.pptx
"""

import sys
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

DPI = 200  # higher = sharper but larger file; 200 is a good balance


def pdf_to_pptx(pdf_path: str, pptx_path: str, dpi: int = DPI) -> None:
    pdf = fitz.open(pdf_path)
    prs = Presentation()

    # Use first page dimensions to set slide size (preserves aspect ratio)
    first = pdf[0].rect
    aspect = first.width / first.height
    slide_h = 6_858_000  # 190.5 mm in EMU (standard 16:9 height)
    slide_w = int(slide_h * aspect)
    prs.slide_width = Emu(slide_w)
    prs.slide_height = Emu(slide_h)

    blank_layout = prs.slide_layouts[6]  # completely blank

    zoom = dpi / 72  # PDF points -> pixels
    for i, page in enumerate(pdf):
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

        # Save to temp bytes buffer
        import io
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(buf, 0, 0, width=Emu(slide_w), height=Emu(slide_h))
        print(f"  Slide {i+1}/{len(pdf)}", end="\r", flush=True)

    pdf.close()
    prs.save(pptx_path)
    print(f"\nSaved: {pptx_path}  ({len(pdf)} slides, DPI={dpi})")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python scripts/pdf_to_pptx.py <input.pdf> <output.pptx>")
        sys.exit(1)
    pdf_to_pptx(sys.argv[1], sys.argv[2])
